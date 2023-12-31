import os
import re
import time
import nltk
import spacy
from PyPDF2 import PdfReader
from docx import Document
from elasticsearch import Elasticsearch
from elasticsearch.helpers import parallel_bulk
from nltk.corpus import stopwords
from nltk.stem import SnowballStemmer
from nltk.tokenize import word_tokenize
from openpyxl import load_workbook
from pptx import Presentation

class DocumentSearcher:
    def __init__(self, data_directory):
        self.config = {
            "data_directory": data_directory,
            "elasticsearch_host": "http://localhost:9200"
        }
        self.es = Elasticsearch(self.config["elasticsearch_host"])
        self.nlp = spacy.load("fr_core_news_sm")
        self._download_stopwords()
        self._download_punkt()

    def _download_stopwords(self):
        # Download French stopwords if not available
        try:
            stopwords.words("french")
        except LookupError:
            nltk.download('stopwords')

    def _download_punkt(self):
        # Download Punkt tokenizer if not available
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt')

    def preprocess_text(self, text):
        # Clean the text by removing HTML tags and non-alphanumeric characters
        clean_text = re.sub('<.*?>', '', text)
        clean_text = re.sub('[^\w\s]', '', clean_text)

        # Lemmatize the text using spaCy
        doc = self.nlp(clean_text)
        lemmatized_tokens = [token.lemma_ if token.lemma_ != "-PRON-" else token.text for token in doc]

        # Tokenize the lemmatized text and perform stemming
        tokens = word_tokenize(" ".join(lemmatized_tokens))
        stemmed_tokens = [SnowballStemmer("french").stem(token) for token in tokens]

        # Filter out stopwords from the stemmed tokens
        filtered_tokens = [token for token in stemmed_tokens if token not in stopwords.words("french")]
        processed_text = " ".join(filtered_tokens)

        return processed_text

    def _get_indexed_files(self):
        # Retrieve the set of already indexed file paths from Elasticsearch
        if not self.es.indices.exists(index="documents"):
            return set()

        query = {
            "query": {"match_all": {}},
            "size": 10000  # Adjust the size to retrieve all indexed documents if necessary
        }

        response = self.es.search(index="documents", body=query)
        hits = response["hits"]["hits"]

        indexed_files = set()

        for hit in hits:
            source = hit["_source"]
            indexed_files.add(source["file_path"])

        return indexed_files

    def collect_data(self):
        # Collect and index the data from the specified directory
        file_paths = self._get_file_paths(self.config["data_directory"])
        indexed_files = self._get_indexed_files()

        documents = self._generate_documents(file_paths, indexed_files)

        start_time = time.time()  # Start measuring the processing time

        self.index_documents(documents)
        end_time = time.time()  # Stop measuring the processing time
        processing_time = end_time - start_time

        if processing_time > 0:
            message = f"Indexing completed. Total time: {processing_time:.2f} seconds."
        else:
            message = "No documents were indexed."

        return message, processing_time

    def _get_file_paths(self, directory):
        # Recursively retrieve all file paths within the specified directory
        file_paths = []

        for root, _, files in os.walk(directory):
            for file in files:
                file_paths.append(os.path.join(root, file))

        return file_paths

    def _extract_text_from_docx(self, file_path):
        # Extract text content from a DOCX file
        doc = Document(file_path)
        text = ""

        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"

        return text.strip()

    def _extract_text_from_xlsx(self, file_path):
        # Extract text content from an XLSX file
        wb = load_workbook(filename=file_path, read_only=True)
        text = ""

        for sheet in wb.sheetnames:
            ws = wb[sheet]

            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        text += str(cell) + "\n"

        return text.strip()

    def _extract_text_from_pptx(self, file_path):
        # Extract text content from a PPTX file
        prs = Presentation(file_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text.strip()

    def _extract_text_from_pdf(self, file_path):
        # Extract text content from a PDF file
        with open(file_path, "rb") as file:
            pdf = PdfReader(file)
            text = ""

            for page in pdf.pages:
                text += page.extract_text()

        return text.strip()

    def _generate_documents(self, file_paths, indexed_files):
        # Generate documents from file paths by extracting their text content
        for file_path in file_paths:
            if file_path not in indexed_files:
                if file_path.endswith(".docx"):
                    document = self._extract_text_from_docx(file_path)
                    yield (file_path, document)
                elif file_path.endswith(".xlsx"):
                    document = self._extract_text_from_xlsx(file_path)
                    yield (file_path, document)
                elif file_path.endswith(".pptx"):
                    document = self._extract_text_from_pptx(file_path)
                    yield (file_path, document)
                elif file_path.endswith(".pdf"):
                    document = self._extract_text_from_pdf(file_path)
                    yield (file_path, document)
                else:
                    print(f"Unsupported file format: {file_path}")
            else:
                print(f"Skipping already indexed file: {file_path}")

    def index_documents(self, documents):
        # Index the extracted documents in Elasticsearch
        actions = [
            {
                "_op_type": "index",
                "_index": "documents",
                "_id": i,
                "_source": {
                    "file_path": file_path,
                    "processed_text": self.preprocess_text(text)
                }
            }
            for i, (file_path, text) in enumerate(documents)
        ]

        for success, info in parallel_bulk(self.es, actions, index="documents"):
            if not success:
                print(f"Failed to index document: {info}")

    def preprocess_query(self, query):
        # Preprocess the query by tokenizing, stemming, and filtering out stopwords
        tokens = word_tokenize(query)
        stemmed_tokens = [SnowballStemmer("french").stem(token) for token in tokens]
        filtered_tokens = [token for token in stemmed_tokens if token not in stopwords.words("french")]
        processed_query = " ".join(filtered_tokens)
        return processed_query

    def search_files(self, query):
        # Search for files matching the processed query in Elasticsearch
        search_query = {
            "query": {
                "match": {
                    "processed_text": query
                }
            },
            "size": 4,  # Adjust the size to retrieve the desired number of results
            "sort": [
                {
                    "_score": {
                        "order": "desc"  # Sort by descending score to get the most relevant results first
                    }
                }
            ]
        }

        response = self.es.search(index="documents", body=search_query)

        hits = response["hits"]["hits"]
        results = []

        for hit in hits:
            source = hit["_source"]
            result = {
                "file_path": source["file_path"],
                "score": hit["_score"]
            }
            results.append(result)

        return results

    def paragraphs_containing_answer(self, file_path, query):
        # Returns the paragraphs, cells, shapes, or pages that contain the answer to the query in the given file, along with the processing time as a tuple of (matching_content, processing_time)

        try:
            processed_query = self.preprocess_query(query)
            processed_query_tokens = processed_query.split()  # Split the processed query into tokens

            if file_path.endswith(".docx"):
                # Search for paragraphs containing the answer in a DOCX file
                doc = Document(file_path)
                matching_paragraphs = []
                max_matching_tokens = 0  # Variable to keep track of the maximum number of matching tokens

                start_time = time.time()  # Start measuring the processing time

                for paragraph in doc.paragraphs:
                    processed_text = self.preprocess_text(paragraph.text.lower())  # Convert processed text to lowercase
                    processed_text_tokens = processed_text.split()  # Split the processed text into tokens

                    # Find the number of query tokens present in the paragraph
                    num_matching_tokens = sum(token in processed_text_tokens for token in processed_query_tokens)

                    if num_matching_tokens > max_matching_tokens:
                        max_matching_tokens = num_matching_tokens
                        matching_paragraphs = [paragraph.text]
                    elif num_matching_tokens == max_matching_tokens:
                        matching_paragraphs.append(paragraph.text)

                end_time = time.time()  # Stop measuring the processing time
                processing_time = end_time - start_time

                matching_paragraphs_str = "\n\n-------------------------\n\n".join(matching_paragraphs)
                return matching_paragraphs_str, processing_time

            elif file_path.endswith(".xlsx"):
                # Search for cells containing the answer in an XLSX file
                wb = load_workbook(filename=file_path, read_only=True)
                matching_cells = []

                start_time = time.time()

                for sheet in wb.sheetnames:
                    ws = wb[sheet]

                    for row in ws.iter_rows(values_only=True):
                        for cell in row:
                            if cell:
                                processed_text = self.preprocess_text(str(cell).lower())
                                processed_text_tokens = processed_text.split()

                                num_matching_tokens = sum(
                                    token in processed_text_tokens for token in processed_query_tokens)

                                if num_matching_tokens > 0:
                                    matching_cells.append(str(cell))

                end_time = time.time()
                processing_time = end_time - start_time

                matching_cells_str = "\n\n-------------------------\n\n".join(matching_cells)
                return matching_cells_str, processing_time

            elif file_path.endswith(".pptx"):
                # Search for shapes containing the answer in a PPTX file
                prs = Presentation(file_path)
                matching_shapes = []

                start_time = time.time()

                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            processed_text = self.preprocess_text(shape.text.lower())
                            processed_text_tokens = processed_text.split()

                            num_matching_tokens = sum(
                                token in processed_text_tokens for token in processed_query_tokens)

                            if num_matching_tokens > 0:
                                matching_shapes.append(shape.text)

                end_time = time.time()
                processing_time = end_time - start_time

                matching_shapes_str = "\n\n-------------------------\n\n".join(matching_shapes)
                return matching_shapes_str, processing_time

            elif file_path.endswith(".pdf"):
                # Search for pages containing the answer in a PDF file
                with open(file_path, "rb") as file:
                    pdf = PdfReader(file)
                    matching_pages = []

                    start_time = time.time()

                    for page in pdf.pages:
                        processed_text = self.preprocess_text(page.extract_text().lower())
                        processed_text_tokens = processed_text.split()

                        num_matching_tokens = sum(
                            token in processed_text_tokens for token in processed_query_tokens)

                        if num_matching_tokens > 0:
                            matching_pages.append(page.extract_text())

                    end_time = time.time()
                    processing_time = end_time - start_time

                    matching_pages_str = "\n\n-------------------------\n\n".join(matching_pages)
                    return matching_pages_str, processing_time

            else:
                return "Unsupported file format.", None

        except Exception as e:
            return str(e), None
